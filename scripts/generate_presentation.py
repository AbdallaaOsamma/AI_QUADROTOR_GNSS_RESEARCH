"""Generate FYP PowerPoint presentation for AI-Augmented Quadrotor GNSS-Denied RL project.

Creates a professional academic-style .pptx with 20 slides covering both semesters.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour Palette ──────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)   # HWU-ish dark navy
ACCENT_BLUE = RGBColor(0x00, 0x6D, 0xAA)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x60, 0x60, 0x60)
LIGHT_GREY = RGBColor(0xF0, 0xF0, 0xF0)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

OUTPUT_PATH = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "..", "docs", "FYP_Presentation_Abdalla_Shoaeb.pptx",
)
OUTPUT_PATH = os.path.normpath(OUTPUT_PATH)


# ── Helper Functions ────────────────────────────────────────────────────────

def _add_background(slide, color=DARK_BLUE):
    """Fill the slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_bar(slide, title_text, subtitle_text=None):
    """Add a coloured title bar at the top of a content slide."""
    # Title bar shape
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.1),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Title text
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)

    if subtitle_text:
        sp = tf.add_paragraph()
        sp.text = subtitle_text
        sp.font.size = Pt(14)
        sp.font.color.rgb = LIGHT_BLUE
        sp.alignment = PP_ALIGN.LEFT

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.1),
        Inches(10), Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()


def _add_bullet_body(slide, bullets, top=1.35, left=0.5, width=9.0, height=5.5,
                     font_size=18, sub_bullets=None):
    """Add bullet points to a slide. sub_bullets is a dict mapping bullet index to list of sub-items."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    sub_bullets = sub_bullets or {}

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(6)
        p.level = 0

        if i in sub_bullets:
            for sub in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.text = sub
                sp.font.size = Pt(font_size - 2)
                sp.font.color.rgb = GREY
                sp.space_after = Pt(4)
                sp.level = 1

    return txBox


def _add_table(slide, headers, rows, top=1.5, left=0.5, width=9.0, row_height=0.4):
    """Add a styled table to a slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    col_width = width / num_cols

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(row_height * num_rows),
    )
    table = table_shape.table

    # Header row
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            bg = LIGHT_GREY if ri % 2 == 0 else WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = BLACK
                p.alignment = PP_ALIGN.CENTER

    return table_shape


def _add_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def _add_section_label(slide, text, top=0.0):
    """Add a small section label in the top-right corner."""
    txBox = slide.shapes.add_textbox(
        Inches(7.5), Inches(top + 0.05),
        Inches(2.3), Inches(0.35),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.RIGHT


# ── Slide Builders ──────────────────────────────────────────────────────────

def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_background(slide, DARK_BLUE)

    # Main title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9.0), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI-Augmented Flight Control for Quadrotor UAV"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Using Reinforcement Learning in GNSS-Denied Environments"
    p2.font.size = Pt(24)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER

    # Separator line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.5), Inches(3.7),
        Inches(5.0), Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Details
    details = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(8.0), Inches(2.5))
    tf2 = details.text_frame
    tf2.word_wrap = True
    for text in [
        "Abdalla Shoaeb  |  H00404752",
        "Supervisor: Marah Saleh",
        "School of Engineering and Physical Sciences",
        "Heriot-Watt University Dubai",
        "B50PR -- Final Year Project  |  March 2026",
    ]:
        if tf2.paragraphs[0].text == "":
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)

    _add_notes(slide, (
        "Welcome everyone. This presentation covers the final year project on "
        "AI-augmented flight control for quadrotor UAVs using reinforcement learning "
        "in GNSS-denied environments. The work spans both semesters and covers "
        "system design, implementation, training, and evaluation."
    ))


def build_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Agenda")
    _add_bullet_body(slide, [
        "1.  Semester 1 Recap: Problem, Literature, Initial Design",
        "2.  System Architecture & Observation Pipeline",
        "3.  Reward Function & Simulated VIO Pipeline",
        "4.  Domain Randomisation & Safety Monitor",
        "5.  Training Infrastructure & Results",
        "6.  Ablation Studies (Reward, Frame Stack, DR)",
        "7.  Software Engineering & Quality Assurance",
        "8.  Limitations, Future Work & Conclusion",
    ], font_size=20)
    _add_notes(slide, (
        "Here is the agenda. We will start with a brief recap of Semester 1 work, "
        "then dive into the full system architecture, training pipeline, "
        "experimental results from ablation studies, and conclude with limitations "
        "and future work."
    ))


def build_sem1_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Problem Statement & Motivation", "Semester 1 Recap")
    _add_section_label(slide, "Semester 1")
    _add_bullet_body(slide, [
        "Quadrotor UAVs are critical for inspection, search-and-rescue, and mapping",
        "Indoor / underground environments lack GNSS (GPS) signals",
        "Traditional PID controllers require accurate state estimates from GPS",
        "Model Predictive Control (MPC) needs dynamics models and is compute-heavy",
        "Reinforcement Learning (RL) can learn navigation policies from experience",
        "Key challenge: bridge the sim-to-real gap without GPS feedback",
    ], font_size=18, sub_bullets={
        1: ["Tunnels, warehouses, urban canyons, collapsed buildings"],
        5: ["Domain randomisation + safety monitor approach"],
    })
    _add_notes(slide, (
        "The motivation is straightforward: indoor and underground environments lack GPS. "
        "Traditional controllers fail without accurate position estimates. "
        "RL offers a data-driven alternative that can learn robust policies from simulation. "
        "The main challenge is making these policies transfer to real hardware."
    ))


def build_sem1_literature(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Literature Review Highlights", "Semester 1 Recap")
    _add_section_label(slide, "Semester 1")
    _add_bullet_body(slide, [
        "PPO (Schulman et al., 2017): stable on-policy RL with clipped surrogate",
        "SAC (Haarnoja et al., 2018): off-policy maximum-entropy, better sample efficiency",
        "AirSim (Shah et al., 2018): Unreal Engine simulation with depth + IMU sensors",
        "Domain Randomisation (Tobin et al., 2017): sensor/physics variation for sim-to-real",
        "VIO / Visual-Inertial Odometry: GPS-free state estimation via camera + IMU",
        "Research Gap: most RL-UAV work uses ground-truth state, not VIO estimates",
    ], font_size=17, sub_bullets={
        5: ["This project fills the gap: RL trained on drift-corrupted VIO, not GPS"],
    })
    _add_notes(slide, (
        "The literature review identified that most RL-based UAV work uses privileged "
        "ground-truth state. Our project addresses this gap by training with simulated VIO "
        "estimates that include realistic drift and bias, enforcing the GNSS-denied constraint."
    ))


def build_sem1_design(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Initial System Design & AirSim Setup", "Semester 1 Recap")
    _add_section_label(slide, "Semester 1")
    _add_bullet_body(slide, [
        "Gymnasium-compatible AirSim wrapper with lockstep simulation",
        "84x84 depth camera, body-frame velocity, 4-frame stack",
        "Configurable reward functions via YAML configs",
        "Altitude-hold via moveByVelocityZBodyFrameAsync (prevents Z drift)",
        "Collision detection with timestamp filtering (prevents stale flags)",
        "Reset with retry loop + 0.5s physics settle delay",
        "Semester 1 deliverable: validated environment, ready for training",
    ], font_size=17)
    _add_notes(slide, (
        "Semester 1 focused on building a reliable simulation wrapper. Key decisions include "
        "lockstep simulation for determinism, altitude hold to simplify the action space, "
        "and collision timestamp filtering to prevent false episode terminations."
    ))


def build_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "System Architecture Overview")

    # Architecture as a text-based diagram using a table
    headers = ["Component", "Description", "Key Detail"]
    rows = [
        ["Observation", "Depth 84x84 + VIO velocity", "4-frame stack, GNSS-denied"],
        ["Policy (PPO)", "NatureCNN + MLP, late fusion", "MultiInputPolicy (SB3)"],
        ["Action Space", "Continuous [-1,1]^3", "vx, vy, yaw_rate"],
        ["Reward", "4 components", "Progress + Collision + Smooth + Drift"],
        ["Safety Monitor", "4-layer envelope", "Deployment-time only"],
        ["VIO Simulator", "Dead-reckoning + Gaussian drift", "Replaces GPS truth"],
        ["Domain Rand.", "4 axes of variation", "Depth, spawn, VIO, optic flow"],
    ]
    _add_table(slide, headers, rows, top=1.4)

    _add_notes(slide, (
        "This table summarises the full system architecture. The observation pipeline "
        "combines depth imagery with VIO-estimated velocity. The PPO policy uses "
        "NatureCNN for the image branch and MLP for velocity, with late fusion. "
        "The safety monitor wraps the policy at deployment time."
    ))


def build_observation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Observation Space")
    _add_bullet_body(slide, [
        "Image: Forward depth camera, 84 x 84 pixels, clipped at 20m, normalised [0, 1]",
        "Velocity: Body-frame [vx, vy, yaw_rate] from simulated VIO (NOT GPS ground-truth)",
        "Frame stacking: 4 consecutive frames via VecFrameStack (temporal context)",
        "Gymnasium Dict space: {\"image\": Box(84,84,4), \"velocity\": Box(3,)}",
        "GNSS explicitly disabled: no position data in observation",
        "Body-frame rotation: AirSim global kinematics rotated via yaw matrix",
    ], font_size=17, sub_bullets={
        1: ["VIO velocity includes Gaussian drift + bias injection"],
        4: ["Agent must rely solely on depth + inertial sensing"],
    })
    _add_notes(slide, (
        "The observation space is designed to mirror what a real GNSS-denied drone would see. "
        "No GPS coordinates anywhere in the observation. The velocity comes from simulated VIO "
        "with intentional drift to force the policy to handle estimation errors."
    ))


def build_reward(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Reward Function Design")

    headers = ["Component", "Formula", "Weight", "Purpose"]
    rows = [
        ["Progress", "w_p * vx_body", "+0.5", "Encourage forward flight"],
        ["Collision", "w_c (terminal)", "-100", "Penalise crashes"],
        ["Smoothness", "w_s * ||a - a_prev||", "-0.1", "Reduce action jerk"],
        ["Drift Penalty", "w_d * ||est - truth||", "-0.05", "Penalise VIO divergence"],
    ]
    _add_table(slide, headers, rows, top=1.4)

    _add_bullet_body(slide, [
        "All weights configurable via configs/rewards/*.yaml",
        "Three reward profiles tested: default, aggressive, cautious",
        "Ablation studies isolate contribution of each component",
    ], top=4.0, font_size=16)

    _add_notes(slide, (
        "The reward function has four components. Progress rewards forward flight, "
        "collision is a large negative terminal penalty, smoothness penalises jerky actions, "
        "and the drift penalty discourages policies that exploit VIO estimation errors. "
        "All weights are configurable and we tested three profiles."
    ))


def build_vio(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Simulated VIO Pipeline")
    _add_bullet_body(slide, [
        "Replaces GPS ground-truth velocity with drift-corrupted estimates",
        "Dead-reckoning model with accumulating Gaussian noise",
        "Bias injection: slow-varying bias term added to velocity readings",
        "Drift magnitude configurable (default: sigma = 0.05 m/s per step)",
        "Why it matters: forces policy to be robust to state estimation error",
        "Real VIO (e.g., VINS-Mono) has similar drift characteristics",
        "Without VIO simulation, the agent \"cheats\" using perfect state",
    ], font_size=17, sub_bullets={
        0: ["Central to the GNSS-denied premise of this project"],
        5: ["Our simplified model captures the essential drift behaviour"],
    })
    _add_notes(slide, (
        "The simulated VIO pipeline is one of the key contributions. Without it, "
        "the agent would be training on perfect ground-truth state, which contradicts "
        "the GNSS-denied premise. Our dead-reckoning model with Gaussian drift "
        "captures the essential characteristics of real VIO systems."
    ))


def build_domain_rand(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Domain Randomisation")

    headers = ["Axis", "Method", "Range"]
    rows = [
        ["Depth Noise", "Gaussian additive noise on depth image", "sigma = 5% of depth"],
        ["Spawn Jitter", "Random position + yaw at episode start", "5m radius, 360 deg yaw"],
        ["VIO Drift", "Random drift rate on velocity estimates", "sigma = 0.03-0.08 m/s"],
        ["Optical Flow", "Noise injection on flow-derived velocity", "sigma = 0.02 m/s"],
    ]
    _add_table(slide, headers, rows, top=1.4)

    _add_bullet_body(slide, [
        "Purpose: improve policy robustness for sim-to-real transfer",
        "Trained policies must generalise across sensor noise distributions",
        "Toggled via configs/train_ppo_dr.yaml",
    ], top=4.2, font_size=16)

    _add_notes(slide, (
        "Domain randomisation is applied across four axes. This exposes the policy "
        "to a range of sensor conditions during training, so it learns representations "
        "that are invariant to specific noise characteristics. This is critical for "
        "eventual sim-to-real transfer."
    ))


def build_safety(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Safety Monitor Architecture")

    headers = ["Layer", "Function", "Trigger", "Action"]
    rows = [
        ["1. Velocity Clamp", "Hard-limit speeds", "v > v_max", "Clamp to bounds"],
        ["2. Proximity Brake", "Scale vx by obstacle dist", "ROI depth < 1.5m", "vx -> 0.2x"],
        ["3. Altitude Guard", "Flag altitude deviation", "delta_z > 1.0m", "Alert + correct"],
        ["4. Emergency Stop", "Zero all commands", "Collision / timeout", "Full stop"],
    ]
    _add_table(slide, headers, rows, top=1.4)

    _add_bullet_body(slide, [
        "Architecturally separated from the RL policy (deployment wrapper only)",
        "Centre 30% ROI for proximity detection (avoids peripheral noise)",
        "Configurable limits via configs/safety.yaml",
        "58 unit tests validate safety logic independently",
    ], top=4.2, font_size=16)

    _add_notes(slide, (
        "The safety monitor is a four-layer envelope that wraps the RL policy at deployment. "
        "It is intentionally separated from the policy so that safety guarantees do not depend "
        "on learned behaviour. Each layer addresses a different failure mode."
    ))


def build_training_infra(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Training Infrastructure")
    _add_bullet_body(slide, [
        "Algorithm: PPO (Stable-Baselines3) with MultiInputPolicy",
        "Simulation: AirSim lockstep (simContinueForTime + simPause)",
        "Parallel envs: SubprocVecEnv for N simultaneous AirSim instances",
        "Training throughput: ~6 FPS (lockstep depth rendering is the bottleneck)",
        "Checkpointing: every 10k steps + best model tracking via eval callback",
        "Monitoring: TensorBoard logging of rewards, episode length, loss terms",
        "ONNX export: trained policy exportable for edge deployment",
        "7 training runs completed, ~300k steps each, across ablation configs",
    ], font_size=17, sub_bullets={
        2: ["DummyVecEnv for N=1, SubprocVecEnv(start_method='spawn') for N>1"],
        6: ["Target: NVIDIA Jetson platform (future work)"],
    })
    _add_notes(slide, (
        "Training uses PPO from Stable-Baselines3 with lockstep AirSim simulation. "
        "We support parallel environments for faster training. The main throughput "
        "bottleneck is depth rendering at 6 FPS. We completed 7 training runs "
        "across different ablation configurations."
    ))


def build_results_reward_ablation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Results: Reward Component Ablation")

    headers = ["Configuration", "Ep. Reward", "Ep. Length", "Key Observation"]
    rows = [
        ["Full Reward (baseline)", "338", "354", "Balanced performance"],
        ["No Smoothness", "417 (+23%)", "389", "Higher reward, more aggressive"],
        ["Progress Only", "361 (+7%)", "264", "Shorter episodes, less cautious"],
    ]
    _add_table(slide, headers, rows, top=1.4, row_height=0.45)

    _add_bullet_body(slide, [
        "Removing smoothness penalty yields +23% reward but more aggressive flying",
        "Progress-only reward produces shorter episodes (faster but riskier)",
        "Full reward provides best balance between speed and safety",
        "All three variants learn effective obstacle avoidance",
    ], top=3.8, font_size=16)

    _add_notes(slide, (
        "The reward ablation shows that smoothness penalty reduces raw reward but produces "
        "smoother, safer flight. The No Smoothness variant gets +23% reward but flies more "
        "aggressively. Progress Only has shorter episodes, suggesting it takes more risks. "
        "The full reward provides the best trade-off for deployment."
    ))


def build_results_frame_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Results: Frame Stack Ablation")

    headers = ["Configuration", "Ep. Reward", "Ep. Length", "Delta vs Baseline"]
    rows = [
        ["1-Frame Stack", "494", "469", "+84% reward (BEST)"],
        ["4-Frame Stack (baseline)", "269", "319", "Baseline"],
    ]
    _add_table(slide, headers, rows, top=1.4, row_height=0.45)

    _add_bullet_body(slide, [
        "Surprising finding: 1-frame stack significantly outperforms 4-frame stack",
        "Hypothesis: with body-frame velocity already in the observation, temporal",
        "  context from frame stacking may be redundant and adds noise",
        "1-frame reduces observation dimensionality by 4x (faster learning)",
        "Implication: VIO velocity may provide sufficient temporal information",
        "Challenges prior assumption that frame stacking always helps",
    ], top=3.4, font_size=16)

    _add_notes(slide, (
        "This is the most surprising result. The 1-frame stack agent achieves 84% higher "
        "reward than the 4-frame baseline. Our hypothesis is that since body-frame velocity "
        "is already in the observation, the temporal context from frame stacking is redundant. "
        "The larger observation space from 4 frames may actually slow learning."
    ))


def build_results_dr(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Results: Domain Randomisation")

    headers = ["Configuration", "Ep. Reward", "Ep. Length", "Note"]
    rows = [
        ["Without DR (baseline)", "338", "354", "Trained on clean data"],
        ["With DR", "291 (-14%)", "326", "Trained with sensor noise"],
    ]
    _add_table(slide, headers, rows, top=1.4, row_height=0.45)

    _add_bullet_body(slide, [
        "DR reduces in-simulation reward by 14% (expected trade-off)",
        "DR-trained policies are designed for better real-world transfer",
        "The reward drop reflects the added difficulty of noisy observations",
        "True benefit of DR is measured at deployment, not in training",
        "Literature consistently shows DR improves sim-to-real transfer",
    ], top=3.4, font_size=16)

    _add_notes(slide, (
        "Domain randomisation reduces training reward by 14%, which is expected -- "
        "the agent is training on harder, noisier data. The real benefit of DR is "
        "measured during deployment on hardware, not during simulation training. "
        "The literature strongly supports DR for sim-to-real transfer."
    ))


def build_results_training_progress(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Results: Training Progress (PPO Base)")

    headers = ["Metric", "Early Training", "Late Training (500k)", "Trend"]
    rows = [
        ["Episode Reward", "197", "217", "+10% improvement"],
        ["Episode Length", "243", "257", "Longer survival"],
        ["Collision Rate", "High", "Decreasing", "Learning to avoid"],
        ["Policy Loss", "Variable", "Converging", "Stabilising"],
    ]
    _add_table(slide, headers, rows, top=1.4, row_height=0.45)

    _add_bullet_body(slide, [
        "PPO baseline shows steady learning over 500k steps",
        "Reward increases by 10%, episode length increases by 6%",
        "Collision rate decreases as training progresses",
        "Additional training steps expected to yield further improvement",
    ], top=3.8, font_size=16)

    _add_notes(slide, (
        "The PPO base run shows clear learning progress over 500k steps. "
        "Reward increases by about 10% and episodes get longer, indicating the agent "
        "is learning to survive and navigate. Collision rates decrease over time. "
        "Longer training runs would likely produce further improvement."
    ))


def build_software_engineering(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Software Engineering & Quality")

    headers = ["Metric", "Value"]
    rows = [
        ["Unit Tests", "58 passing (AirSim-free)"],
        ["Lint Errors", "0 (ruff)"],
        ["YAML Configs", "14 config files"],
        ["Training Profiles", "4 (base, fast, DR, multi-env)"],
        ["Reward Profiles", "3 (default, aggressive, cautious)"],
        ["Ablation Configs", "4 (no_smooth, progress, frame, DR)"],
        ["ONNX Export", "Automated script"],
    ]
    _add_table(slide, headers, rows, top=1.4, row_height=0.38)

    _add_bullet_body(slide, [
        "All tests run without AirSim (mock at boundary)",
        "Modular architecture: environments / training / evaluation / safety / deployment",
        "Config-driven: all hyperparameters in YAML, never hardcoded",
    ], top=5.0, font_size=15)

    _add_notes(slide, (
        "The project follows strong software engineering practices. "
        "58 tests all pass without AirSim by mocking at the boundary. "
        "Zero lint errors with ruff. All hyperparameters are in YAML configs. "
        "The modular architecture separates concerns cleanly."
    ))


def build_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Known Limitations")
    _add_bullet_body(slide, [
        "VIO model is simplified (dead-reckoning + Gaussian noise)",
        "  - Real VIO (VINS-Mono, ORB-SLAM3) has more complex failure modes",
        "SAC comparison not yet completed (PPO-only results so far)",
        "Statistical analysis (ANOVA, t-tests) not finalised across all runs",
        "Physics domain randomisation limited (mass, inertia, wind not yet varied)",
        "Hardware prototype is planned future work, not current scope",
        "  - Project scope refined to simulation-based research with ONNX export",
        "Training throughput limited to ~6 FPS by depth rendering",
    ], font_size=16)
    _add_notes(slide, (
        "Key limitations: our VIO model is simplified compared to real VIO systems. "
        "SAC comparison is not yet complete. Statistical tests are not finalised. "
        "The hardware prototype has been scoped as future work -- this project "
        "focuses on simulation-based research with ONNX export for future deployment."
    ))


def build_future_work(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Future Work")
    _add_bullet_body(slide, [
        "SAC implementation and comparative evaluation vs PPO",
        "Real VIO integration (VINS-Mono or ORB-SLAM3) replacing simplified model",
        "Full physics domain randomisation (mass, inertia, wind, motor gains)",
        "Curriculum learning: progressively harder environments",
        "Hardware validation on Quanser QDrone 2 (Jetson Xavier NX + RealSense)",
        "  - ONNX inference via ROS2 bridge to MAVLink flight controller",
        "Optuna-driven hyperparameter optimisation",
        "Multi-waypoint navigation with goal-conditioned policies",
    ], font_size=17, sub_bullets={
        4: ["Planned future phase, not current project scope"],
    })
    _add_notes(slide, (
        "Future work includes implementing SAC for comparison, integrating a real VIO "
        "pipeline, expanding domain randomisation to physics parameters, and eventually "
        "deploying on the QDrone 2 hardware platform. Curriculum learning and "
        "Optuna-based hyperparameter tuning are also planned."
    ))


def build_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, WHITE)
    _add_title_bar(slide, "Conclusion")
    _add_bullet_body(slide, [
        "Designed and implemented a complete RL training framework for GNSS-denied UAV navigation",
        "Enforced GNSS-denied constraint via simulated VIO (not post-hoc GPS removal)",
        "Trained PPO agent successfully learns obstacle avoidance from depth + VIO velocity",
        "Ablation studies reveal 1-frame stack outperforms 4-frame by 84% (novel finding)",
        "4-layer safety monitor provides deployment-ready safety guarantees",
        "Domain randomisation prepares policies for sim-to-real transfer",
        "58 tests, 14 configs, modular architecture -- production-quality codebase",
        "ONNX export pipeline established for future hardware deployment",
    ], font_size=17)

    # Add a final thank you
    txBox = slide.shapes.add_textbox(Inches(2.0), Inches(6.2), Inches(6.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank you -- Questions?"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    _add_notes(slide, (
        "In conclusion, this project delivers a complete RL framework for GNSS-denied "
        "quadrotor navigation. The key contributions are: enforcing GNSS-denied constraints "
        "via simulated VIO, a novel finding that 1-frame stacking outperforms 4-frame, "
        "a 4-layer safety monitor, and a production-quality codebase. "
        "Thank you for listening. I'm happy to take questions."
    ))


# ── Main ────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Build all 20 slides in order
    build_title_slide(prs)          # 1
    build_agenda(prs)               # 2
    build_sem1_problem(prs)         # 3
    build_sem1_literature(prs)      # 4
    build_sem1_design(prs)          # 5
    build_architecture(prs)         # 6
    build_observation(prs)          # 7
    build_reward(prs)               # 8
    build_vio(prs)                  # 9
    build_domain_rand(prs)          # 10
    build_safety(prs)               # 11
    build_training_infra(prs)       # 12
    build_results_reward_ablation(prs)   # 13
    build_results_frame_stack(prs)       # 14
    build_results_dr(prs)                # 15
    build_results_training_progress(prs) # 16
    build_software_engineering(prs)      # 17
    build_limitations(prs)               # 18
    build_future_work(prs)               # 19
    build_conclusion(prs)                # 20

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
